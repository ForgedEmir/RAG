"""Create realistic B2B PowerPoint files in data/sample/ for testing."""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt

SAMPLE = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "sample"))


def create_proposition_commerciale():
    """Sales deck: 5 slides, text + pricing table."""
    prs = Presentation()

    # Slide 1 - Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Proposition Commerciale - Oracle LoreKeeper"
    s.placeholders[1].text = "Solution RAG Enterprise\nQ2 2026"

    # Slide 2 - Context
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Contexte et Enjeux"
    s.placeholders[1].text = (
        "Les entreprises font face a 3 defis majeurs :\n"
        "1. Documents eparpilles sur 5+ plateformes\n"
        "2. 68% du temps perdu a chercher de l'information\n"
        "3. Risque de non-conformite RGPD\n\n"
        "Notre solution RAG centralise et securise l'acces a la connaissance."
    )

    # Slide 3 - Pricing table
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "Grille Tarifaire"
    tb.text_frame.paragraphs[0].font.size = Pt(28)

    tbl = s.shapes.add_table(5, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    for ci, h in enumerate(["Offre", "Utilisateurs", "Stockage", "Prix/mois"]):
        tbl.cell(0, ci).text = h
    for ri, row in enumerate([
        ["Starter", "1-10", "5 Go", "99 EUR"],
        ["Business", "11-50", "50 Go", "399 EUR"],
        ["Enterprise", "51-200", "500 Go", "999 EUR"],
        ["Custom", "200+", "Illimite", "Sur devis"],
    ]):
        for ci, val in enumerate(row):
            tbl.cell(ri + 1, ci).text = val

    # Slide 4 - ROI
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Retour sur Investissement"
    s.placeholders[1].text = (
        "Gains mesures chez nos clients :\n"
        "- Temps de recherche reduit de 73%\n"
        "- Productivite equipe +35%\n"
        "- Conformite RGPD : 100% des documents audites\n"
        "- ROI moyen : 4.2x en 12 mois\n\n"
        "Etude de cas : Cabinet Dupont & Associes\n"
        "Avant : 2h/jour de recherche par consultant\n"
        "Apres : 15 min/jour avec Oracle LoreKeeper"
    )

    # Slide 5 - Next steps
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Prochaines Etapes"
    s.placeholders[1].text = (
        "1. POC gratuit de 30 jours\n"
        "2. Integration avec votre SSO (Azure AD / Okta)\n"
        "3. Migration des documents existants\n"
        "4. Formation equipe (2h)\n"
        "5. Go-live et support 24/7"
    )

    path = os.path.join(SAMPLE, "proposition_commerciale_2026.pptx")
    prs.save(path)
    print(f"Created: proposition_commerciale_2026.pptx (5 slides, pricing table)")


def create_rapport_financier():
    """Financial quarterly report: 4 slides, dense tables with numbers."""
    prs = Presentation()

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Rapport Financier T1 2026"
    s.placeholders[1].text = "Direction Financiere - Confidentiel"

    # Revenue by region
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "Chiffre d'Affaires par Region"
    tb.text_frame.paragraphs[0].font.size = Pt(24)

    tbl = s.shapes.add_table(6, 4, Inches(0.5), Inches(1.3), Inches(9), Inches(3.5)).table
    for ci, h in enumerate(["Region", "T1 2025", "T1 2026", "Variation"]):
        tbl.cell(0, ci).text = h
    for ri, row in enumerate([
        ["France", "1 250 000 EUR", "1 480 000 EUR", "+18.4%"],
        ["DACH (DE/AT/CH)", "890 000 EUR", "1 020 000 EUR", "+14.6%"],
        ["Benelux", "420 000 EUR", "510 000 EUR", "+21.4%"],
        ["UK", "680 000 EUR", "590 000 EUR", "-13.2%"],
        ["TOTAL", "3 240 000 EUR", "3 600 000 EUR", "+11.1%"],
    ]):
        for ci, val in enumerate(row):
            tbl.cell(ri + 1, ci).text = val

    # Expenses
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "Ventilation des Charges"
    tb.text_frame.paragraphs[0].font.size = Pt(24)

    tbl = s.shapes.add_table(7, 3, Inches(1), Inches(1.3), Inches(8), Inches(4)).table
    for ci, h in enumerate(["Poste", "Montant", "% du CA"]):
        tbl.cell(0, ci).text = h
    for ri, row in enumerate([
        ["Masse salariale", "1 800 000 EUR", "50.0%"],
        ["Infrastructure cloud", "360 000 EUR", "10.0%"],
        ["Marketing & Commercial", "540 000 EUR", "15.0%"],
        ["R&D", "450 000 EUR", "12.5%"],
        ["Frais generaux", "180 000 EUR", "5.0%"],
        ["TOTAL CHARGES", "3 330 000 EUR", "92.5%"],
    ]):
        for ci, val in enumerate(row):
            tbl.cell(ri + 1, ci).text = val

    # KPIs
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "KPIs Cles"
    s.placeholders[1].text = (
        "MRR : 300 000 EUR (+22% YoY)\n"
        "ARR : 3 600 000 EUR\n"
        "Clients actifs : 47 (+12 vs T4 2025)\n"
        "Churn : 2.1% (objectif < 5%)\n"
        "NPS : 72 (excellent)\n"
        "CAC : 8 500 EUR (objectif < 10 000 EUR)\n"
        "LTV/CAC : 5.2x"
    )

    path = os.path.join(SAMPLE, "rapport_financier_t1_2026.pptx")
    prs.save(path)
    print(f"Created: rapport_financier_t1_2026.pptx (4 slides, financial tables)")


def create_architecture_technique():
    """Technical architecture deck: 4 slides, specs + formats table."""
    prs = Presentation()

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Architecture Technique - Oracle LoreKeeper"
    s.placeholders[1].text = "Documentation technique v3.2"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Stack Technique"
    s.placeholders[1].text = (
        "Backend : FastAPI + Python 3.11\n"
        "Vector DB : Qdrant (cloud)\n"
        "Embeddings : BGE-M3 via FastEmbed (ONNX)\n"
        "LLM : Cerebras (primary) + Groq (fallback)\n"
        "Cache : Redis (semantic cache + sessions)\n"
        "Monitoring : Langfuse + Sentry\n"
        "Frontend : React 18 + Vite\n"
        "Auth : JWT + SSO (Azure AD, Okta)\n"
        "Deploiement : Docker + Kubernetes"
    )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Pipeline RAG"
    s.placeholders[1].text = (
        "Ingestion :\n"
        "  1. Upload fichier (PDF, DOCX, XLSX, PPTX, CSV, Email)\n"
        "  2. Parsing multi-format avec detection encoding\n"
        "  3. Chunking table-aware (1200 chars, 200 overlap)\n"
        "  4. Embedding BGE-M3 (dense + sparse)\n"
        "  5. Indexation Qdrant + BM25 corpus\n\n"
        "Retrieval :\n"
        "  1. Query expansion (multi-query LLM)\n"
        "  2. Hybrid search (vector + BM25 + RRF k=60)\n"
        "  3. Cross-encoder reranking (Jina v2)\n"
        "  4. HyDE fallback si confiance < 0.005"
    )

    # Formats table
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "Formats Supportes"
    tb.text_frame.paragraphs[0].font.size = Pt(24)

    tbl = s.shapes.add_table(8, 3, Inches(1), Inches(1.3), Inches(8), Inches(4)).table
    for ci, h in enumerate(["Format", "Parser", "Specificites"]):
        tbl.cell(0, ci).text = h
    for ri, row in enumerate([
        [".pdf", "LlamaParse / pypdf / OCR", "Tableaux, scan, multi-colonnes"],
        [".docx", "python-docx", "Headings, headers/footers, tables"],
        [".xlsx", "openpyxl", "Cellules fusionnees, multi-feuilles"],
        [".xls", "xlrd", "Format legacy Excel 97-2003"],
        [".pptx", "python-pptx", "Slides, tableaux, notes"],
        [".csv", "csv (stdlib)", "Auto-detect delimiteur et encoding"],
        [".eml/.msg", "email/extract-msg", "Headers, body, pieces jointes"],
    ]):
        for ci, val in enumerate(row):
            tbl.cell(ri + 1, ci).text = val

    path = os.path.join(SAMPLE, "architecture_technique_v3.pptx")
    prs.save(path)
    print(f"Created: architecture_technique_v3.pptx (4 slides, tech specs + formats table)")


if __name__ == "__main__":
    create_proposition_commerciale()
    create_rapport_financier()
    create_architecture_technique()
    print(f"\nAll files created in {SAMPLE}")
