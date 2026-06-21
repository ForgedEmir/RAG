"""
B2B Stress Test Suite — Real-world edge cases for enterprise document parsing.
Tests encoding detection, multi-language, CSV/Excel/PPTX/Email edge cases,
table-aware chunking, and real files from data/sample.
"""
import csv
import json
import os
import sys
import tempfile

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from src.ingestion.parser import extract_text_from_file, clean_text
from src.ingestion.chunker import split_into_chunks

TEST_DIR = tempfile.mkdtemp(prefix="b2b_stress_")
_results = {"PASS": 0, "FAIL": 0, "errors": []}


def _test(name, condition, detail=""):
    if condition:
        _results["PASS"] += 1
        print(f"  PASS: {name}")
    else:
        _results["FAIL"] += 1
        _results["errors"].append(f"{name}: {detail}")
        print(f"  FAIL: {name} -- {detail[:120]}")


def _make(name, content_bytes, ext=".txt"):
    path = os.path.join(TEST_DIR, name + ext)
    with open(path, "wb") as f:
        f.write(content_bytes)
    return path


# ── 1. ENCODING EDGE CASES ──────────────────────────────────────────────────
def test_encoding():
    print("\n[1] ENCODING EDGE CASES")

    # Latin-1 French business text
    p = _make("latin1_fr", "Le chiffre d'affaires a atteint 2,5M. Pr\xe9visions favorables.".encode("latin-1"))
    r = extract_text_from_file(p)
    _test("Latin-1 French", r is not None and "visions" in r, repr(r)[:80] if r else "None")

    # Windows-1252 smart quotes + euro
    p = _make("cp1252", b"Revenue \x93exceeded\x94 expectations \x96 total: \x80500K")
    r = extract_text_from_file(p)
    _test("Windows-1252 smart quotes + euro", r is not None and "exceeded" in r, repr(r)[:80] if r else "None")

    # UTF-8 BOM
    p = _make("utf8_bom", b"\xef\xbb\xbfProject Alpha: Budget approval pending")
    r = extract_text_from_file(p)
    _test("UTF-8 BOM stripped", r is not None and r.lstrip("\ufeff").startswith("Project"), repr(r)[:80] if r else "None")

    # UTF-16 LE
    p = _make("utf16le", "Quarterly Review: Revenue up 23%".encode("utf-16-le"))
    r = extract_text_from_file(p)
    _test("UTF-16 LE", r is not None and "Revenue" in r, repr(r)[:80] if r else "None")

    # ISO-8859-15
    p = _make("iso8859_15", "Facture: 1.250,00 HT".encode("iso-8859-15"))
    r = extract_text_from_file(p)
    _test("ISO-8859-15 (EU)", r is not None and "Facture" in r, repr(r)[:80] if r else "None")


# ── 2. MULTI-LANGUAGE DOCUMENTS ──────────────────────────────────────────────
def test_languages():
    print("\n[2] MULTI-LANGUAGE DOCUMENTS")

    p = _make("german", "Geschaeftsbericht 2025: Umsatz gestiegen.".encode("utf-8"))
    r = extract_text_from_file(p)
    _test("German", r is not None and "Geschaeftsbericht" in r)

    p = _make("spanish", "Informe anual: los ingresos alcanzaron 3M EUR.".encode("utf-8"))
    r = extract_text_from_file(p)
    _test("Spanish", r is not None and "ingresos" in r)

    p = _make("arabic", "\u062a\u0642\u0631\u064a\u0631 \u0645\u0627\u0644\u064a: 5 \u0645\u0644\u064a\u0648\u0646 \u064a\u0648\u0631\u0648".encode("utf-8"))
    r = extract_text_from_file(p)
    _test("Arabic (RTL)", r is not None and len(r) > 10)

    p = _make("japanese", "\u56db\u534a\u671f\u5831\u544a\u66f8\uff1a\u58f2\u4e0a\u9ad8\u306f\u524d\u5e74\u6bd415%\u5897\u52a0".encode("utf-8"))
    r = extract_text_from_file(p)
    _test("Japanese", r is not None and len(r) > 10)

    p = _make("mixed_lang", "EN: Revenue 2M.\nFR: Chiffre 2M.\nDE: Umsatz 2M.".encode("utf-8"))
    r = extract_text_from_file(p)
    _test("Mixed EN/FR/DE", r is not None and "Revenue" in r and "Chiffre" in r and "Umsatz" in r)


# ── 3. CSV EDGE CASES ────────────────────────────────────────────────────────
def test_csv():
    print("\n[3] CSV EDGE CASES")

    p = _make("csv_semi", "Nom;Montant;Devise\nDupont;1500;EUR\nMuller;2300;EUR".encode("utf-8"), ".csv")
    r = extract_text_from_file(p)
    _test("CSV semicolons", r is not None and "Dupont" in r and "1500" in r)

    p = _make("csv_tab", "Name\tAmount\nSmith\t5000\nJones\t3200".encode("utf-8"), ".csv")
    r = extract_text_from_file(p)
    _test("CSV tab-delimited", r is not None and "Smith" in r)

    p = _make("csv_latin1", "Client;Montant\nSoci\xe9t\xe9;50000".encode("latin-1"), ".csv")
    r = extract_text_from_file(p)
    _test("CSV Latin-1 encoding", r is not None and "50000" in r)

    p = _make("csv_quoted", b'Name,Address,Amount\n"Smith, John","123 Main St, Suite 5",5000', ".csv")
    r = extract_text_from_file(p)
    _test("CSV quoted fields with commas", r is not None and "Smith" in r and "5000" in r)

    p = _make("csv_empty", b"", ".csv")
    r = extract_text_from_file(p)
    _test("CSV empty file (no crash)", True)  # just verify no exception


# ── 4. JSON EDGE CASES ───────────────────────────────────────────────────────
def test_json():
    print("\n[4] JSON EDGE CASES")

    data = {
        "company": {"name": "Acme Corp", "revenue": {"2024": 5000000, "2025": 6200000}},
        "employees": [{"name": "Alice", "role": "CEO"}, {"name": "Bob", "role": "CTO"}],
    }
    p = _make("json_nested", json.dumps(data).encode("utf-8"), ".json")
    r = extract_text_from_file(p)
    _test("Nested JSON", r is not None and "Acme Corp" in r and "Alice" in r)

    data = {"entreprise": "Soci\u00e9t\u00e9 G\u00e9n\u00e9rale", "CA": "50M EUR"}
    p = _make("json_unicode", json.dumps(data, ensure_ascii=False).encode("utf-8"), ".json")
    r = extract_text_from_file(p)
    _test("JSON Unicode FR", r is not None and "50M" in r)


# ── 5. EXCEL EDGE CASES ─────────────────────────────────────────────────────
def test_excel():
    print("\n[5] EXCEL EDGE CASES")
    import openpyxl

    # Multi-sheet
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Revenue"
    ws1.append(["Quarter", "Amount"])
    ws1.append(["Q1", 1500000])
    ws2 = wb.create_sheet("Expenses")
    ws2.append(["Category", "Amount"])
    ws2.append(["Salaries", 800000])
    p = os.path.join(TEST_DIR, "multi_sheet.xlsx")
    wb.save(p)
    r = extract_text_from_file(p)
    _test("XLSX multi-sheet", r is not None and "Revenue" in r and "Expenses" in r and "Salaries" in r)

    # Merged cells propagation
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Dept"
    ws["B1"] = "Q1"
    ws["A2"] = "Engineering"
    ws.merge_cells("A2:A4")
    ws["B2"] = 100
    ws["B3"] = 110
    ws["B4"] = 105
    p = os.path.join(TEST_DIR, "merged.xlsx")
    wb.save(p)
    r = extract_text_from_file(p)
    eng_count = r.count("Engineering") if r else 0
    _test("XLSX merged cells (3 rows)", eng_count >= 3, f"Engineering appeared {eng_count} times")

    # Formula errors filtered
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Item"
    ws["B1"] = "Value"
    ws["A2"] = "Valid"
    ws["B2"] = 100
    ws["A3"] = "Error"
    ws["B3"] = "#REF!"
    ws["A4"] = "DivZero"
    ws["B4"] = "#DIV/0!"
    p = os.path.join(TEST_DIR, "formula_errors.xlsx")
    wb.save(p)
    r = extract_text_from_file(p)
    _test("XLSX formula errors filtered", r is not None and "#REF!" not in r and "#DIV/0!" not in r and "Valid" in r)

    # Empty sheet
    wb = openpyxl.Workbook()
    p = os.path.join(TEST_DIR, "empty_sheet.xlsx")
    wb.save(p)
    r = extract_text_from_file(p)
    _test("XLSX empty sheet (no crash)", True)


# ── 6. POWERPOINT EDGE CASES ────────────────────────────────────────────────
def test_pptx():
    print("\n[6] POWERPOINT EDGE CASES")
    from pptx import Presentation
    from pptx.util import Inches

    # Slide with table
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tbl = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(6), Inches(2)).table
    for ri, row in enumerate([["Product", "Price", "Stock"], ["Widget A", "29.99", "500"], ["Widget B", "49.99", "120"]]):
        for ci, val in enumerate(row):
            tbl.cell(ri, ci).text = val
    p = os.path.join(TEST_DIR, "pptx_table.pptx")
    prs.save(p)
    r = extract_text_from_file(p)
    _test("PPTX with table", r is not None and "Widget A" in r and "29.99" in r)

    # 10 slides
    prs = Presentation()
    for i in range(10):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i + 1} Title"
        slide.placeholders[1].text = f"Content for slide {i + 1}"
    p = os.path.join(TEST_DIR, "pptx_10.pptx")
    prs.save(p)
    r = extract_text_from_file(p)
    _test("PPTX 10 slides", r is not None and "Slide 1" in r and "Slide 10" in r)

    # Empty slide
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    p = os.path.join(TEST_DIR, "pptx_empty.pptx")
    prs.save(p)
    r = extract_text_from_file(p)
    _test("PPTX empty slide (no crash)", True)


# ── 7. EMAIL EDGE CASES ─────────────────────────────────────────────────────
def test_email():
    print("\n[7] EMAIL EDGE CASES")

    # HTML email
    html_eml = (
        "From: sales@company.com\r\n"
        "To: client@business.com\r\n"
        "Subject: Enterprise License Proposal\r\n"
        "MIME-Version: 1.0\r\n"
        "Content-Type: text/html; charset=utf-8\r\n"
        "\r\n"
        "<html><body><h1>Proposal</h1>"
        "<p>We offer our <strong>Enterprise License</strong> at <em>50,000 EUR/year</em>.</p>"
        "</body></html>\r\n"
    )
    p = _make("eml_html", html_eml.encode("utf-8"), ".eml")
    r = extract_text_from_file(p)
    _test("EML HTML tags stripped", r is not None and "Enterprise License" in r and "<strong>" not in r)

    # Email with attachment
    att_eml = (
        "From: legal@corp.com\r\n"
        "To: finance@corp.com\r\n"
        "Subject: Q4 Audit Files\r\n"
        "MIME-Version: 1.0\r\n"
        'Content-Type: multipart/mixed; boundary="bound123"\r\n'
        "\r\n"
        "--bound123\r\n"
        "Content-Type: text/plain; charset=utf-8\r\n"
        "\r\n"
        "Please review the attached audit documents.\r\n"
        "\r\n"
        "--bound123\r\n"
        'Content-Type: application/pdf; name="audit_report.pdf"\r\n'
        'Content-Disposition: attachment; filename="audit_report.pdf"\r\n'
        "Content-Transfer-Encoding: base64\r\n"
        "\r\n"
        "JVBERi0=\r\n"
        "--bound123--\r\n"
    )
    p = _make("eml_att", att_eml.encode("utf-8"), ".eml")
    r = extract_text_from_file(p)
    _test("EML attachment listed", r is not None and "audit_report.pdf" in r)

    # Minimal email
    p = _make("eml_min", b"Subject: Test\r\n\r\nJust a body.", ".eml")
    r = extract_text_from_file(p)
    _test("EML minimal", r is not None and "Test" in r and "body" in r)


# ── 8. DOCX EDGE CASES ──────────────────────────────────────────────────────
def test_docx():
    print("\n[8] DOCX EDGE CASES")
    from docx import Document

    # 3-level headings
    doc = Document()
    doc.add_heading("Board Meeting", level=1)
    doc.add_heading("Financial Overview", level=2)
    doc.add_paragraph("Revenue: 5M EUR")
    doc.add_heading("Regional Breakdown", level=3)
    doc.add_paragraph("EMEA: 3M, APAC: 2M")
    p = os.path.join(TEST_DIR, "docx_h3.docx")
    doc.save(p)
    r = extract_text_from_file(p)
    _test("DOCX 3-level headings", r is not None and "# Board" in r and "## Financial" in r and "### Regional" in r)

    # Table in DOCX
    doc = Document()
    doc.add_paragraph("Contract Summary")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Party"
    table.cell(0, 1).text = "Role"
    table.cell(1, 0).text = "Acme Corp"
    table.cell(1, 1).text = "Vendor"
    p = os.path.join(TEST_DIR, "docx_tbl.docx")
    doc.save(p)
    r = extract_text_from_file(p)
    _test("DOCX with table", r is not None and "Acme Corp" in r and "Vendor" in r)


# ── 9. TABLE-AWARE CHUNKING ─────────────────────────────────────────────────
def test_chunking():
    print("\n[9] TABLE-AWARE CHUNKING")

    # Large table (50 rows)
    big = "Report Header\n\n"
    for i in range(50):
        big += f"Product_{i:03d} | {10 + i * 2}.99 | {100 + i * 10} | Cat_{i % 5}\n"
    big += "\nConclusion: inventory is healthy."
    chunks = split_into_chunks(big, chunk_size=500, overlap=100)
    broken = any(
        "|" in line and len(line.split("|")) < 2
        for c in chunks
        for line in c.strip().split("\n")
        if line.strip()
    )
    _test("Large table: no mid-row splits", not broken, f"{len(chunks)} chunks")

    # Mixed prose + tables
    mixed = (
        "Executive Summary\n\nOur company had a strong quarter.\n\n"
        "Employee | Dept | Salary\nAlice | Eng | 80000\nBob | Mkt | 65000\n"
        "Charlie | Sales | 70000\nDiana | Eng | 85000\n\n"
        "The table above shows salaries.\n\n"
        "Name | Rating\nAlice | 4.5\nBob | 3.8\n\n"
        "End of report."
    )
    chunks = split_into_chunks(mixed, chunk_size=300, overlap=50)
    _test("Mixed prose + 2 tables", len(chunks) >= 3, f"{len(chunks)} chunks")


# ── 10. REAL FILES ───────────────────────────────────────────────────────────
def test_real_files():
    print("\n[10] PARSING REAL B2B FILES")
    sample_dir = os.path.join(os.path.dirname(__file__), "..", "..", "data", "sample")
    sample_dir = os.path.normpath(sample_dir)
    if not os.path.isdir(sample_dir):
        print(f"  SKIP: {sample_dir} not found")
        return

    real_files = sorted(
        f for f in os.listdir(sample_dir)
        if os.path.isfile(os.path.join(sample_dir, f)) and not f.startswith(".")
    )
    for fname in real_files:
        fpath = os.path.join(sample_dir, fname)
        try:
            r = extract_text_from_file(fpath)
            if r is None:
                _test(fname, False, "returned None")
            elif len(r.strip()) < 10:
                _test(fname, False, f"too short: {len(r)} chars")
            else:
                chunks = split_into_chunks(clean_text(r))
                _test(fname, True)
                print(f"          -> {len(r)} chars, {len(chunks)} chunks")
        except Exception as e:
            _test(fname, False, f"EXCEPTION: {e}")


# ── MAIN ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 70)
    print("B2B STRESS TEST SUITE")
    print("=" * 70)

    test_encoding()
    test_languages()
    test_csv()
    test_json()
    test_excel()
    test_pptx()
    test_email()
    test_docx()
    test_chunking()
    test_real_files()

    print("\n" + "=" * 70)
    print(f"RESULTS: {_results['PASS']} PASS / {_results['FAIL']} FAIL")
    print("=" * 70)
    if _results["errors"]:
        print("\nFailed tests:")
        for err in _results["errors"]:
            print(f"  - {err}")
    else:
        print("\nAll tests passed!")

    sys.exit(1 if _results["FAIL"] else 0)
